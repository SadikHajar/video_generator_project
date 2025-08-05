import os
import json
from datetime import datetime


from dotenv import load_dotenv
import google.generativeai as genai
from fpdf import FPDF
from pptx import Presentation
from pptx.util import Inches, Pt
from create_video_from_script import create_video_from_script
import re
from langdetect import detect
from image_manager import ImageManager

load_dotenv()

class ScriptGenerator:
    def __init__(self):
        print("🔍 Chargement de la clé Gemini...")
        self.api_key = os.getenv("GEMINI_API_KEY")

        if not self.api_key:
            raise Exception("❌ GEMINI_API_KEY non trouvée dans .env")

        genai.configure(api_key=self.api_key)
        self.model = genai.GenerativeModel("gemini-1.5-flash")
        # Initialiser le gestionnaire d'images
        self.image_manager = ImageManager()

    def generate_training_script(self, user_prompt):
        # Détection de la langue de l'utilisateur
        lang = detect(user_prompt)
        print(f"🌐 Langue détectée: {lang}")

        if lang == 'en':
            system_prompt = """
You are an expert in professional training content creation.
Given a training request, generate a full and structured script.

The script must contain a MINIMUM of 8-12 scenes for a complete and in-depth training. Each scene should include:
- A catchy and descriptive title
- A detailed voice-over script (180-250 words)
- A descriptive keyword for visual elements (we will use this to search for relevant images automatically)
- 3-4 key takeaway points

IMPORTANT: For visual elements, provide ONLY descriptive keywords that describe the scene content (example: "artificial intelligence", "neural network", "programming code", "data analysis", etc.). Do NOT provide URLs - we will handle image retrieval automatically.

Recommended structure:
1. Introduction and context
2. Basics and fundamentals (2-3 scenes)
3. Intermediate concepts (3-4 scenes)
4. Advanced concepts (2-3 scenes)
5. Practical applications (1-2 scenes)
6. Conclusion and perspectives

Respond ONLY with a valid JSON in this exact format:
{
  "titre_formation": "Full and catchy training title",
  "description": "Detailed training description",
  "objectifs": ["Objective 1", "Objective 2", "Objective 3", "Objective 4"],
  "scenes": [
    {
      "numero": 1,
      "titre": "Detailed scene title",
      "voix_off": "Full and detailed voice-over text (180-250 words)",
      "elements_visuels": "descriptive keywords for image search",
      "points_cles": ["Key point 1", "Key point 2", "Key point 3"]
    }
  ]
}
"""
        else:
            system_prompt = """
Tu es un expert en création de contenu de formation professionnelle.
À partir d'une demande de formation, tu dois générer un script complet et structuré.

Le script doit contenir MINIMUM 8-12 scènes pour une formation complète et approfondie. Chaque scène doit avoir :
- Un titre accrocheur et descriptif
- Un script de voix off détaillé (180-250 mots)
- Des mots-clés descriptifs pour les éléments visuels (nous utiliserons ceux-ci pour rechercher automatiquement des images pertinentes)
- Des points clés à retenir (3-4 points par scène)

IMPORTANT: Pour les éléments visuels, fournis UNIQUEMENT des mots-clés descriptifs qui décrivent le contenu de la scène (exemple: "intelligence artificielle", "réseau de neurones", "code de programmation", "analyse de données", etc.). Ne fournis PAS d'URLs - nous nous occupons automatiquement de la récupération d'images.

Structure recommandée:
1. Introduction et contexte
2. Bases et fondamentaux (2-3 scènes)
3. Concepts intermédiaires (3-4 scènes)
4. Concepts avancés (2-3 scènes)
5. Applications pratiques (1-2 scènes)
6. Conclusion et perspectives

Réponds UNIQUEMENT avec un JSON valide dans ce format exact:
{
  "titre_formation": "Titre complet et accrocheur",
  "description": "Description détaillée de la formation",
  "objectifs": ["Objectif 1", "Objectif 2", "Objectif 3", "Objectif 4"],
  "scenes": [
    {
      "numero": 1,
      "titre": "Titre scène détaillé",
      "voix_off": "Texte voix off complet et détaillé (180-250 mots)",
      "elements_visuels": "mots-clés descriptifs pour recherche d'image",
      "points_cles": ["Point clé 1", "Point clé 2", "Point clé 3"]
    }
  ]
}
"""
        print(f"\n🤖 Génération du script pour: {user_prompt}")
        
        # Combiner le system prompt avec la demande utilisateur
        full_prompt = f"{system_prompt}\n\nDemande de formation: {user_prompt}"
        
        response = self.model.generate_content(
            full_prompt,
            generation_config={"temperature": 0.7}
        )

        content = response.text.strip()
        
        # DEBUG : Afficher la réponse brute pour voir ce qui est retourné
        print("🔍 DEBUG - Réponse de l'API:")
        print("=" * 50)
        print(content)  # Afficher TOUTE la réponse pour voir ce qui se passe
        print("=" * 50)
        print(f"📏 Longueur de la réponse: {len(content)} caractères")
        
        # Essayer d'extraire le JSON de la réponse
        json_content = self._extract_json_from_response(content)
        
        print("🔍 DEBUG - JSON extrait:")
        print("=" * 30)
        print(json_content[:300] + "..." if len(json_content) > 300 else json_content)
        print("=" * 30)

        try:
            script_data = json.loads(json_content)
        except json.JSONDecodeError as e:
            print(f"❌  : {e}")
            print(f"🔍 Contenu à parser: {json_content[:200]}...")
            raise Exception(f"Erreur de parsing JSON: {e}")
        
        # Valider et corriger les URLs d'images avec le gestionnaire d'images
        script_data = self.image_manager.validate_and_fix_image_urls(script_data)

        self._validate_script_structure(script_data)
        print("✅ Script généré avec succès!")
        return script_data

    def _extract_json_from_response(self, content):
        """Extrait le JSON de la réponse de l'API, même s'il y a du texte avant/après"""
        import re
        
        # Enlever les marqueurs de code markdown
        content = re.sub(r'```json\s*', '', content)
        content = re.sub(r'```\s*$', '', content)
        content = content.strip()
        
        # Chercher un objet JSON complet dans la réponse
        json_start = content.find('{')
        if json_start != -1:
            # Compter les accolades pour trouver la fin du JSON
            brace_count = 0
            json_end = json_start
            
            for i, char in enumerate(content[json_start:], json_start):
                if char == '{':
                    brace_count += 1
                elif char == '}':
                    brace_count -= 1
                    if brace_count == 0:
                        json_end = i + 1
                        break
            
            if brace_count == 0:
                json_content = content[json_start:json_end]
                return json_content
        
        # Si pas de JSON complet trouvé, essayer ligne par ligne
        lines = content.split('\n')
        cleaned_lines = []
        in_json = False
        brace_count = 0
        
        for line in lines:
            line = line.strip()
            if line.startswith('{') and not in_json:
                in_json = True
                brace_count = 0
            
            if in_json:
                cleaned_lines.append(line)
                # Compter les accolades dans cette ligne
                brace_count += line.count('{') - line.count('}')
                
                # Si on a fermé toutes les accolades, on a fini
                if brace_count <= 0 and line.endswith('}'):
                    break
        
        if cleaned_lines:
            return '\n'.join(cleaned_lines)
        
        # Dernier recours : retourner le contenu tel quel
        return content
    

    def _validate_script_structure(self, script_data):
        required_fields = ['titre_formation', 'description', 'scenes']
        for field in required_fields:
            if field not in script_data:
                raise ValueError(f"Champ manquant: {field}")

        if not isinstance(script_data['scenes'], list):
            raise ValueError("Les scènes doivent être une liste")

        if len(script_data['scenes']) < 8:
            print(f"⚠️ Attention: Seulement {len(script_data['scenes'])} scènes générées. Recommandé: 8-12 scènes minimum.")

        for i, scene in enumerate(script_data['scenes']):
            for field in ['numero', 'titre', 'voix_off', 'elements_visuels']:
                if field not in scene:
                    raise ValueError(f"Scène {i+1} - champ manquant: {field}")
              # Valider que elements_visuels est une URL valide
            if not scene['elements_visuels'].startswith(('http://', 'https://')):
                print(f"⚠️ Scène {i+1}: L'élément visuel ne semble pas être une URL valide")    

    def save_script_to_file(self, script_data, filename=None, folder="data"):
        if not os.path.exists(folder):
            os.makedirs(folder)

        if filename is None:
            timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
            filename = f"script_formation_{timestamp}.json"

        full_path = os.path.join(folder, filename)

        with open(full_path, 'w', encoding='utf-8') as f:
            json.dump(script_data, f, ensure_ascii=False, indent=2)

        print(f"💾 Script JSON sauvegardé dans: {full_path}")
        return full_path

    def json_to_pdf(self, script_data, folder="script"):
        if not os.path.exists(folder):
            os.makedirs(folder)

        timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
        pdf_path = os.path.join(folder, f"script_formation_{timestamp}.pdf")

        pdf = FPDF()
        pdf.set_auto_page_break(auto=True, margin=15)
        pdf.add_page()
        pdf.set_font("Arial", 'B', 16)

        pdf.cell(0, 10, script_data.get('titre_formation', 'Formation'), ln=True)

        pdf.set_font("Arial", '', 12)
        pdf.multi_cell(0, 10, script_data.get('description', ''))

        pdf.ln(5)
        pdf.set_font("Arial", 'I', 12)
        pdf.cell(0, 10, f"Durée estimée: {script_data.get('duree_estimee', 'N/A')}")
        pdf.ln()
        pdf.cell(0, 10, f"Niveau: {script_data.get('niveau', 'N/A')}")
        pdf.ln(10)

        # Objectifs
        pdf.set_font("Arial", 'B', 14)
        pdf.cell(0, 10, "Objectifs :", ln=True)
        pdf.set_font("Arial", '', 12)
        for obj in script_data.get('objectifs', []):
            pdf.cell(0, 8, f"- {obj}", ln=True)
        pdf.ln(10)

        # Scènes
        for scene in script_data.get('scenes', []):
            pdf.set_font("Arial", 'B', 13)
            pdf.cell(0, 10, f"Scène {scene['numero']}: {scene['titre']}", ln=True)

            pdf.set_font("Arial", '', 12)
            pdf.multi_cell(0, 10, f"Voix off:\n{scene['voix_off']}")
            pdf.ln(3)
            pdf.multi_cell(0, 10, f"Éléments visuels:\n{scene['elements_visuels']}")
            pdf.ln(3)

            if scene.get('points_cles'):
                pdf.set_font("Arial", 'I', 12)
                pdf.cell(0, 8, "Points clés:", ln=True)
                pdf.set_font("Arial", '', 12)
                for point in scene['points_cles']:
                    pdf.cell(0, 8, f"- {point}", ln=True)
            pdf.ln(10)

        pdf.output(pdf_path)
        print(f"📄 PDF sauvegardé dans: {pdf_path}")
        return pdf_path

    def display_script_summary(self, script_data):
        print("\n" + "="*60)
        print(f"📚 FORMATION: {script_data['titre_formation']}")
        print("="*60)
        print(f"📝 Description: {script_data['description']}")

        if 'duree_estimee' in script_data:
            print(f"⏱️  Durée estimée: {script_data['duree_estimee']}")

        if 'niveau' in script_data:
            print(f"📊 Niveau: {script_data['niveau']}")

        if 'objectifs' in script_data:
            print(f"🎯 Objectifs:")
            for obj in script_data['objectifs']:
                print(f"   • {obj}")

        print(f"\n🎬 SCÈNES ({len(script_data['scenes'])}):")
        for scene in script_data['scenes']:
            print(f"\n{scene['numero']}. {scene['titre']}")
            print(f"   📝 Voix off: {len(scene['voix_off'])} caractères")
            print(f"   🖼️  Visuels: {scene['elements_visuels'][:100]}...")
            if 'points_cles' in scene:
                print(f"   🔑 Points clés: {', '.join(scene['points_cles'])}")
               

    def json_to_ppt(self, script_data, folder="script"):
        if not os.path.exists(folder):
            os.makedirs(folder)

        timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
        ppt_path = os.path.join(folder, f"script_formation_{timestamp}.pptx")

        prs = Presentation()
        
        # Slide de titre
        title_slide = prs.slides.add_slide(prs.slide_layouts[0])
        title = title_slide.shapes.title
        subtitle = title_slide.placeholders[1]
        
        title.text = script_data.get('titre_formation', 'Formation')
        subtitle.text = f"{script_data.get('description', '')}\n\nDurée: {script_data.get('duree_estimee', 'N/A')} | Niveau: {script_data.get('niveau', 'N/A')}"

        # Slide des objectifs
        if script_data.get('objectifs'):
            obj_slide = prs.slides.add_slide(prs.slide_layouts[1])
            obj_slide.shapes.title.text = "Objectifs de la formation"
            content = obj_slide.placeholders[1]
            
            obj_text = ""
            for obj in script_data['objectifs']:
                obj_text += f"• {obj}\n"
            content.text = obj_text

        # Slides pour chaque scène
        for scene in script_data.get('scenes', []):
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            
            # Titre de la scène
            slide.shapes.title.text = f"Scène {scene['numero']}: {scene['titre']}"
            
            # Contenu
            content = slide.placeholders[1]
            scene_text = f"Voix off:\n{scene['voix_off']}\n\n"
            scene_text += f"Éléments visuels:\n{scene['elements_visuels']}\n\n"
            
            if scene.get('points_cles'):
                scene_text += "Points clés:\n"
                for point in scene['points_cles']:
                    scene_text += f"• {point}\n"
                    
            content.text = scene_text

        prs.save(ppt_path)
        print(f"📊 PowerPoint sauvegardé dans: {ppt_path}")
        return ppt_path
    
